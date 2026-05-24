from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional
from xml.etree import ElementTree
from zipfile import ZipFile

from fastapi import HTTPException

from app.services.pdf.pdf_loader import extract_pdf_text


@dataclass
class PageText:
    page: int
    text: str
    used_ocr: bool = False


def extract_document_text(path: Path) -> List[PageText]:
    extension = path.suffix.lower()
    if extension == ".pdf":
        return extract_pdf_text(path)
    if extension == ".docx":
        return extract_docx_text(path)
    if extension == ".pptx":
        return extract_pptx_text(path)
    raise HTTPException(status_code=400, detail=f"Unsupported document type: {extension or 'unknown'}")


def extract_docx_text(path: Path) -> List[PageText]:
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="Word document support is not installed.") from exc

    document = Document(path)
    parts: List[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
        if any(break_tag.get(qn("w:type")) == "page" for run in paragraph.runs for break_tag in run._element.findall(qn("w:br"))):
            parts.append("__DOCUMIND_PAGE_BREAK__")

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return paginate_text_blocks(parts, target_pages=read_docx_page_count(path))


def read_docx_page_count(path: Path) -> Optional[int]:
    try:
        with ZipFile(path) as archive:
            xml = archive.read("docProps/app.xml")
        root = ElementTree.fromstring(xml)
    except Exception:
        return None

    for child in root:
        if child.tag.endswith("Pages") and child.text and child.text.isdigit():
            pages = int(child.text)
            return pages if pages > 0 else None
    return None


def paginate_text_blocks(parts: Iterable[str], max_chars: int = 1800, target_pages: Optional[int] = None) -> List[PageText]:
    blocks = [part.strip() for part in parts if part.strip()]
    if target_pages and target_pages > 0:
        return paginate_text_blocks_to_target(blocks, target_pages)

    return paginate_text_blocks_by_size(blocks, max_chars)


def paginate_text_blocks_to_target(parts: List[str], target_pages: int) -> List[PageText]:
    best_pages: List[PageText] = []
    best_score: tuple[int, int] | None = None

    for max_chars in range(500, 2501, 25):
        pages = paginate_text_blocks_by_size(parts, max_chars)
        score = (abs(len(pages) - target_pages), abs(max_chars - 1000))
        if best_score is None or score < best_score:
            best_pages = pages
            best_score = score
        if len(pages) == target_pages:
            return pages

    return best_pages


def paginate_text_blocks_by_size(parts: Iterable[str], max_chars: int) -> List[PageText]:
    pages: List[PageText] = []
    current: List[str] = []
    current_length = 0

    def flush() -> None:
        nonlocal current, current_length
        text = "\n".join(current).strip()
        if text:
            pages.append(PageText(page=len(pages) + 1, text=text))
        current = []
        current_length = 0

    for part in parts:
        block = part.strip()
        if not block:
            continue
        if block == "__DOCUMIND_PAGE_BREAK__":
            flush()
            continue
        separator_length = 1 if current else 0
        projected_length = current_length + len(block) + separator_length
        if current and projected_length > max_chars:
            flush()
            separator_length = 0
        current.append(block)
        current_length += len(block) + separator_length

    flush()
    return pages


def extract_pptx_text(path: Path) -> List[PageText]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise HTTPException(status_code=500, detail="PowerPoint document support is not installed.") from exc

    presentation = Presentation(path)
    pages: List[PageText] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))

        text = "\n".join(parts).strip()
        pages.append(PageText(page=index, text=text))

    return pages
